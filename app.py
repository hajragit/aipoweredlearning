from flask import Flask, render_template, request, jsonify
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.chat_models import ChatOpenAI
from langchain.schema import Document
from langchain.chains import RetrievalQA
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
import openai
import requests
import json
from flask import Flask, request, render_template, jsonify, send_file
import openai
from pptx import Presentation
from io import BytesIO
import tempfile
from pptx.util import Pt


# Initialize Flask app
app = Flask(__name__)

# Global variables to store vector store and chain
vector_store = None
qa_chain = None

# OpenAI API keys (make sure to handle securely in production)
openai.api_key = ""  # Replace with your actual API key

@app.route("/")
def home():
    return render_template("index.html")

# Route for the understanding_handouts page (where user interacts with the chatbot)
@app.route("/understanding_handouts", methods=["GET", "POST"])
def understanding_handouts():
    global vector_store, qa_chain

    if request.method == "POST":
        # Handle PDF upload
        file = request.files.get("file")
        if not file or not file.filename.endswith(".pdf"):
            return jsonify({"error": "Please upload a valid PDF file."}), 400

        # Process the PDF file
        try:
            pdf_reader = PdfReader(file)
            pdf_text = "".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])

            if pdf_text.strip():
                # Split the text into chunks
                text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
                text_chunks = text_splitter.split_text(pdf_text)

                # Generate embeddings and create the vector store
                embeddings = OpenAIEmbeddings(openai_api_key=openai.api_key)
                documents = [Document(page_content=chunk) for chunk in text_chunks]
                vector_store = FAISS.from_documents(documents, embeddings)

                # Create the QA chain for querying the document
                llm = ChatOpenAI(openai_api_key=openai.api_key, temperature=0)
                qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=vector_store.as_retriever())

                return jsonify({"message": "PDF uploaded and processed successfully!"}), 200
            else:
                return jsonify({"error": "Could not extract text from the PDF."}), 400

        except Exception as e:
            return jsonify({"error": f"An error occurred while processing the PDF: {str(e)}"}), 500

    return render_template("understanding_handouts.html")
@app.route("/query", methods=["POST"])
def query():
    global qa_chain

    if not qa_chain:
        return jsonify({"error": "No document uploaded. Please upload a document first."}), 400

    user_query = request.json.get("query", "").strip()
    if user_query:
        try:
            response = qa_chain.run(user_query)
            return jsonify({"response": response}), 200
        except Exception as e:
            return jsonify({"error": f"An error occurred while processing the query: {str(e)}"}), 500

    return jsonify({"error": "Query cannot be empty."}), 400

# Define the OpenAI API keys
# Define the OpenAI API keys

#openai.api_key = ""




def load_pdf(file):
    pdf = PdfReader(file)
    text = ""
    for page in pdf.pages:
        text += page.extract_text()
    return text


# Route for quiz generation
@app.route("/quiz-generation", methods=["GET", "POST"])
def quiz_generation():
    if request.method == "POST":
        # Check if a file is uploaded
        file = request.files.get("file")
        topic = request.form.get("topic")

        if file:
            # Check if the file is a PDF
            if not file.filename.endswith(".pdf"):
                return jsonify({"error": "Please upload a valid PDF file."}), 400

            try:
                # Extract text from the PDF (you need to define this function)
                pdf_text = load_pdf(file)
                
                if pdf_text.strip():
                    # Split the text into chunks
                    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
                    text_chunks = text_splitter.split_text(pdf_text)

                    # Generate embeddings and create the vector store
                    documents = [Document(page_content=chunk) for chunk in text_chunks]
                    embeddings = OpenAIEmbeddings(openai_api_key=openai.api_key)
                    vector_store = FAISS.from_documents(documents, embeddings)

                    # Create the QA chain for querying the document
                    llm = ChatOpenAI(openai_api_key=openai.api_key, temperature=0)
                    qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=vector_store.as_retriever())

                    # Generate the quiz using the query
                    query = "Based on the document, create 5 multiple-choice questions (MCQs) with 4 options each. Clearly indicate the correct answer."
                    
                    # Run the query to get the quiz
                    result = qa_chain.run(query)
                    return jsonify({"quiz": result}), 200

                else:
                    return jsonify({"error": "Could not extract text from the PDF."}), 400

            except Exception as e:
                return jsonify({"error": f"An error occurred while processing the PDF: {str(e)}"}), 500

        elif topic:
            try:
                # Define the API URL and headers
                url = "https://api.openai.com/v1/chat/completions"
                headers = {
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {openai.api_key}"
                }

                # Define the messages for the chat prompt
                messages = [
                    {"role": "system", "content": "You are a helpful assistant."},
                    {"role": "user", "content": f"Based on the following topic, create 5 multiple-choice questions (MCQs) with 4 options each. Indicate the correct answer clearly.\n\n{topic}"}
                ]

                # Define the data to send in the POST request
                data = {
                    "model": "gpt-3.5-turbo",
                    "messages": messages
                }

                # Send the POST request to OpenAI API
                response = requests.post(url, headers=headers, data=json.dumps(data))

                # Check if the response is successful
                if response.status_code == 200:
                    response_data = response.json()
                    quiz_content = response_data['choices'][0]['message']['content'].strip()
                    return jsonify({"quiz": quiz_content}), 200
                else:
                    return jsonify({"error": f"Request failed with status code {response.status_code}"}), 500

            except Exception as e:
                return jsonify({"error": f"An error occurred: {str(e)}"}), 500

        else:
            return jsonify({"error": "Please upload a document or provide a topic."}), 400

    return render_template("quiz.html")

# Set your OpenAI API key
#openai.api_key = ""

def generate_slide_content(prompt, model="gpt-4o"):
    """
    Generate content for slides using GPT-4 with the updated API.
    """
    
    
    try:
        response = openai.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "You are an expert at creating presentation slides."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=300
        )
        content = response.choices[0].message.content
        print(f"GPT-4 Response:\n{content}\n")
        return content.strip()
    except Exception as e:
        print(f"Error generating content: {e}")
        return "Content generation failed."



def create_presentation(slide_data, file_name="presentation.pptx"):
    """
    Create a PowerPoint presentation using the provided slide data.
    """
    presentation = Presentation()

    for slide_title, slide_content in slide_data.items():
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Set the slide title
        title.text = slide_title

        # Format the content with bullet points
        content.text = ""
        for bullet in slide_content.split("\n"):
            if bullet.strip():
                content.text += f"{bullet.strip()}\n"
        # Adjust font size for all paragraphs
        for paragraph in content.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(17)  # Set the font size to 13pt

    presentation.save(file_name)
    print(f"Presentation saved as {file_name}")
    # Save the presentation to a BytesIO object instead of a file
    pptx_file = BytesIO()
    presentation.save(pptx_file)
    pptx_file.seek(0)  # Move the pointer to the beginning of the file
    return pptx_file

@app.route("/slides", methods=["GET", "POST"])
def lecture_slides():
    if request.method == "POST":
        topic = request.form.get("topic")

        # Validate input
        if not topic or len(topic.strip()) == 0:
            return jsonify({"error": "Topic cannot be empty."}), 400

        # Generate slide outline
        outline_prompt = f"""
Generate an outline for a PowerPoint presentation on the topic '{topic}'.
Please follow these guidelines:
1. Provide 5 slides, each with a clear and concise title (not too long).
2. Under each slide title, list 3-5 bullet points that are focused, relevant, and brief. Avoid long paragraphs.
3. Ensure each slide covers an important aspect of the topic, with each bullet point presenting one key idea or concept.
4. Do not include any filler or irrelevant content.
"""
        outline = generate_slide_content(outline_prompt)

        # Parse the outline into slide titles
        slide_titles = [
            line.strip("-* ").strip() for line in outline.split("\n") if line.strip().startswith(("-", "*"))
        ]

        if not slide_titles:
            return jsonify({"error": "Could not generate slide titles. Please refine the topic and try again."}), 400

        # Generate content for each slide using the updated concise prompt
        slide_data = {}
        for title in slide_titles:
            slide_prompt = f"Create a concise, bullet-point outline for a slide titled: '{title}'. Include key ideas in 3-5 short bullet points."
            slide_content = generate_slide_content(slide_prompt)
            slide_data[title] = slide_content if slide_content else "Content unavailable."

        # Create the PowerPoint presentation
        pptx_file = create_presentation(slide_data, file_name=f"{topic.replace(' ', '_')}.pptx")

        # Send the file for download
        return send_file(pptx_file, download_name=f"{topic.replace(' ', '_')}.pptx", as_attachment=True)

    return render_template("lecture_slides.html")

# Run the Flask app
if __name__ == "__main__":
    app.run(debug=True)